# Sangeeth Deleep Menon
# NUID: 002524579
# CS5330 Final Project — Generate PPTX from presentation content
# Spring 2026

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colours (matching the HTML theme) ──────────────────────────────────────
BLACK  = RGBColor(0x0a, 0x0a, 0x0a)
CARD   = RGBColor(0x16, 0x16, 0x17)
WHITE  = RGBColor(0xf5, 0xf5, 0xf7)
GRAY   = RGBColor(0x86, 0x86, 0x8b)
BLUE   = RGBColor(0x29, 0x97, 0xff)
GREEN  = RGBColor(0x30, 0xd1, 0x58)
YELLOW = RGBColor(0xff, 0xd6, 0x0a)
RED    = RGBColor(0xff, 0x45, 0x3a)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

BASE = os.path.dirname(os.path.abspath(__file__))
FIGURES = os.path.join(BASE, '..', 'Figures')
ROOT    = os.path.join(BASE, '..')

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=BLACK):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height,
        text='', fontsize=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, bg_color=None, border_color=None,
        italic=False, wrap=True):
    """Add a text box; returns the shape."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    if border_color:
        txBox.line.color.rgb = border_color
        txBox.line.width = Pt(1)
    return txBox


def eyebrow(slide, text, top=0.45):
    box(slide, 0, top, 13.33, 0.4, text=text,
        fontsize=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def title(slide, text, top=0.85, fontsize=36):
    box(slide, 0.5, top, 12.33, 1.0, text=text,
        fontsize=fontsize, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def subtitle(slide, text, top=1.9, fontsize=18):
    box(slide, 1.0, top, 11.33, 0.6, text=text,
        fontsize=fontsize, color=GRAY, align=PP_ALIGN.CENTER)


def divider(slide, top=1.8):
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(6.17), Inches(top), Inches(1.0), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def card_box(slide, left, top, width, height,
             label='', value='', desc='',
             label_color=GRAY, value_color=WHITE,
             bg=CARD, border=RGBColor(0x2d, 0x2d, 0x2f),
             value_size=28, desc_size=12):
    """Draw a card with label / big value / description."""
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = border
    shape.line.width = Pt(1)

    if label:
        box(slide, left + 0.15, top + 0.15, width - 0.3, 0.25,
            text=label, fontsize=9, bold=True, color=label_color)
    if value:
        box(slide, left + 0.15, top + 0.4, width - 0.3, 0.7,
            text=value, fontsize=value_size, bold=True, color=value_color)
    if desc:
        box(slide, left + 0.15, top + 1.05, width - 0.3, 0.5,
            text=desc, fontsize=desc_size, color=GRAY)


def add_image(slide, path, left, top, width, height=None):
    if not os.path.exists(path):
        return
    if height:
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    else:
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 Inches(width))


# ── SLIDE 1 — Title ─────────────────────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 0.5, 13.33, 0.4,
    text='CS5330  ·  FINAL PROJECT  ·  SPRING 2026',
    fontsize=10, color=GRAY, align=PP_ALIGN.CENTER)
box(s, 0.5, 1.1, 12.33, 1.6,
    text='Real-Time Object Recognition\nand AR Overlay',
    fontsize=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, top=2.8)
box(s, 1.0, 3.0, 11.33, 0.6,
    text='Using deep learning to turn any everyday object into an augmented reality anchor',
    fontsize=16, color=GRAY, align=PP_ALIGN.CENTER)
box(s, 0, 4.2, 13.33, 0.5,
    text='Sangeeth Deleep Menon  ·  NUID 002524579\nKhoury College of Computer Sciences  ·  Northeastern University',
    fontsize=13, color=GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 2 — The Problem ────────────────────────────────────────────────────
s = add_slide(); bg(s)
eyebrow(s, 'THE PROBLEM')
box(s, 0.5, 0.85, 12.33, 1.2,
    text='Traditional AR needs a printed marker.\nWhat if the object itself was the marker?',
    fontsize=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
card_box(s, 1.0, 2.3, 5.0, 1.5,
         label='PROJECT 4 GAVE US', value='AR rendering',
         desc='3D graphics overlaid using a chessboard calibration target',
         value_size=22)
card_box(s, 6.5, 2.3, 5.0, 1.5,
         label='PROJECT 5 GAVE US', value='Deep learning',
         desc='CNNs and Vision Transformers trained to classify images',
         value_size=22)
# Blue highlight bar
shape = s.shapes.add_shape(1, Inches(1.0), Inches(4.1), Inches(11.33), Inches(0.75))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x0a, 0x19, 0x29)
shape.line.color.rgb = RGBColor(0x1a, 0x6f, 0xc4); shape.line.width = Pt(1)
box(s, 1.0, 4.15, 11.33, 0.6,
    text='This project combines both — no marker required',
    fontsize=17, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# ── SLIDE 3 — System Overview ────────────────────────────────────────────────
s = add_slide(); bg(s)
eyebrow(s, 'SYSTEM OVERVIEW')
title(s, 'A three-step pipeline')
steps = [
    ('1', 'Classify', 'Crop a center ROI from the webcam frame and run it through MobileNetV2 to get a class and confidence score'),
    ('2', 'Estimate pose', 'Use a fixed-depth pinhole camera model to back-project the ROI bounding box into 3D space at 0.5 m'),
    ('3', 'Render', 'Project 8 box corners back to pixels with cv2.projectPoints and draw a 3D wireframe box with a floating label tag'),
]
for i, (num, step_title, desc) in enumerate(steps):
    top = 2.1 + i * 1.3
    circ = s.shapes.add_shape(9, Inches(2.0), Inches(top + 0.05), Inches(0.38), Inches(0.38))
    circ.fill.solid(); circ.fill.fore_color.rgb = BLUE
    circ.line.fill.background()
    box(s, 2.0, top + 0.02, 0.38, 0.38, text=num,
        fontsize=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    box(s, 2.55, top, 8.5, 0.35, text=step_title,
        fontsize=16, bold=True, color=WHITE)
    box(s, 2.55, top + 0.35, 8.5, 0.55, text=desc,
        fontsize=13, color=GRAY)

# ── SLIDE 4 — Dataset ────────────────────────────────────────────────────────
s = add_slide(); bg(s)
eyebrow(s, 'DATASET')
title(s, '10 classes  ·  1419 images')

# Table header
rows = [
    ('book',           '161', 'Webcam', GREEN),
    ('cup',            '173', 'Webcam', GREEN),
    ('keyboard',       '152', 'Webcam', GREEN),
    ('pen',            '171', 'Webcam', GREEN),
    ('phone',          '159', 'Webcam', GREEN),
    ('glasses',        '130', 'Web',    GRAY),
    ('headphones',     '110', 'Web',    GRAY),
    ('laptop',         '123', 'Web',    GRAY),
    ('ps5_controller', '102', 'Web',    GRAY),
    ('tablet',         '138', 'Web',    GRAY),
]
col_lefts = [0.5, 4.0, 5.3]
box(s, col_lefts[0], 1.95, 3.0, 0.3, text='CLASS',   fontsize=9, bold=True, color=GRAY)
box(s, col_lefts[1], 1.95, 1.0, 0.3, text='IMAGES',  fontsize=9, bold=True, color=GRAY)
box(s, col_lefts[2], 1.95, 1.5, 0.3, text='SOURCE',  fontsize=9, bold=True, color=GRAY)
for i, (cls, count, src, src_color) in enumerate(rows):
    row_top = 2.3 + i * 0.42
    box(s, col_lefts[0], row_top, 3.0, 0.38, text=cls,   fontsize=13, color=WHITE)
    box(s, col_lefts[1], row_top, 1.0, 0.38, text=count, fontsize=13, color=WHITE)
    box(s, col_lefts[2], row_top, 1.5, 0.38, text=src,   fontsize=12, bold=True, color=src_color)

# Right side cards
card_box(s, 7.2, 2.0, 5.5, 1.3,
         label='TRAIN / VAL / TEST SPLIT',
         value='70%  /  15%  /  15%',
         desc='Fixed seed 42 — fully reproducible',
         value_size=18)
card_box(s, 7.2, 3.45, 5.5, 1.1,
         label='AUGMENTATION (TRAIN ONLY)',
         value='Flip · Jitter · ±15° rotation',
         desc='',
         value_size=14, value_color=WHITE)
card_box(s, 7.2, 4.7, 5.5, 1.15,
         label='KEY OBSERVATION',
         value='',
         desc='Webcam classes outperform web-only classes — data quality matters more than quantity',
         bg=RGBColor(0x0a, 0x19, 0x29),
         border=RGBColor(0x1a, 0x6f, 0xc4),
         desc_size=13)

# ── SLIDE 5 — Architectures ──────────────────────────────────────────────────
s = add_slide(); bg(s)
eyebrow(s, 'MODEL ARCHITECTURES')
title(s, 'Three approaches, one winner')

arch = [
    ('Custom CNN', 'ObjectCNN',
     ['3-layer conv network', '64×64 RGB input', 'BatchNorm + Dropout 0.4', 'Trained from scratch', '~1.2M parameters'],
     '~100% on 5 classes', False),
    ('Vision Transformer', 'ObjectViT',
     ['8×8 patch embedding', '64 patches per image', '4 transformer layers', 'Trained from scratch', '128-dim embeddings'],
     'Slower to converge', False),
    ('Transfer Learning  ✦  Deployed', 'MobileNetV2',
     ['ImageNet pretrained', '224×224 RGB input', 'Fine-tuned classifier head', 'Dropout 0.3', 'Best on noisy web data'],
     '84.1% on 10 classes', True),
]
for i, (tag, name, bullets, result, selected) in enumerate(arch):
    left = 0.4 + i * 4.3
    bg_col  = RGBColor(0x0a, 0x19, 0x29) if selected else CARD
    bdr_col = BLUE if selected else RGBColor(0x2d, 0x2d, 0x2f)
    shape = s.shapes.add_shape(1, Inches(left), Inches(2.0), Inches(4.0), Inches(4.8))
    shape.fill.solid(); shape.fill.fore_color.rgb = bg_col
    shape.line.color.rgb = bdr_col; shape.line.width = Pt(1.5 if selected else 1)
    box(s, left+0.15, 2.1, 3.7, 0.3, text=tag.upper(),
        fontsize=9, bold=True, color=BLUE if selected else GRAY)
    box(s, left+0.15, 2.4, 3.7, 0.45, text=name,
        fontsize=22, bold=True, color=WHITE)
    for j, b in enumerate(bullets):
        box(s, left+0.3, 2.95 + j*0.4, 3.5, 0.35, text='• ' + b,
            fontsize=12, color=GRAY)
    res_bg  = RGBColor(0x0d, 0x2a, 0x45) if selected else RGBColor(0x1a, 0x1a, 0x1c)
    res_col = BLUE if selected else GRAY
    shape2 = s.shapes.add_shape(1, Inches(left+0.15), Inches(6.1), Inches(3.7), Inches(0.45))
    shape2.fill.solid(); shape2.fill.fore_color.rgb = res_bg
    shape2.line.fill.background()
    box(s, left+0.15, 6.1, 3.7, 0.45, text=result,
        fontsize=13, bold=True, color=res_col, align=PP_ALIGN.CENTER)

# ── SLIDE 6 — Results ────────────────────────────────────────────────────────
s = add_slide(); bg(s)
eyebrow(s, 'RESULTS')
title(s, 'Numbers that matter')

card_box(s, 0.4, 1.9, 3.8, 1.6,
         label='TEST ACCURACY', value='84.1%',
         desc='MobileNetV2 on 10-class held-out test set',
         value_color=GREEN, value_size=36,
         bg=RGBColor(0x0a, 0x1a, 0x10), border=RGBColor(0x1a, 0x4a, 0x2a))
card_box(s, 4.57, 1.9, 3.8, 1.6,
         label='LIVE FRAME RATE', value='30 fps',
         desc='Apple M-series with MPS acceleration',
         value_color=BLUE, value_size=36,
         bg=RGBColor(0x0a, 0x19, 0x29), border=RGBColor(0x1a, 0x6f, 0xc4))
card_box(s, 8.73, 1.9, 3.8, 1.6,
         label='STRONGEST CLASS', value='Phone',
         desc='29 of 30 test images correct',
         value_color=WHITE, value_size=36)

img_path = os.path.join(FIGURES, 'MOBILENET_Loss&Accuracy.png')
add_image(s, img_path, left=0.4, top=3.7, width=12.5)

# ── SLIDE 7 — Confusion & Failures ──────────────────────────────────────────
s = add_slide(); bg(s)
eyebrow(s, 'WHERE IT STRUGGLES')
title(s, 'Confusion tells the story')

img_path = os.path.join(FIGURES, 'MOBILENET_ConfusionMatrix.png')
add_image(s, img_path, left=0.4, top=1.85, width=6.0)

card_box(s, 6.9, 1.85, 5.9, 1.4,
         label='GLASSES — WEAKEST CLASS', label_color=RED,
         value='', desc='Web crawl mixed eyeglasses with drinking glasses. The model had no way to separate them.',
         desc_size=13)
card_box(s, 6.9, 3.4, 5.9, 1.4,
         label='TABLET ↔ PHONE CONFUSION', label_color=YELLOW,
         value='', desc='Flat rectangles look alike at certain angles. The AR screenshot shows this exact case live.',
         desc_size=13)
card_box(s, 6.9, 4.95, 5.9, 1.3,
         label='ROOT CAUSE', value='',
         desc='All 5 weak classes used web images only. Webcam-captured classes had near-perfect scores.',
         bg=RGBColor(0x0a, 0x19, 0x29), border=RGBColor(0x1a, 0x6f, 0xc4), desc_size=13)

# ── SLIDE 8 — Live Demo ──────────────────────────────────────────────────────
s = add_slide(); bg(s)
eyebrow(s, 'LIVE DEMO')
title(s, 'AR overlay in action')

img_path = os.path.join(ROOT, 'ar_screenshot.png')
add_image(s, img_path, left=0.4, top=1.85, width=6.0)

box(s, 6.9, 1.85, 6.0, 0.7,
    text='A tablet is in frame. The model predicts phone at 84% confidence.',
    fontsize=15, color=WHITE)
box(s, 6.9, 2.65, 6.0, 0.8,
    text='The 3D wireframe box and label are rendered correctly around the detected region even though the class is wrong — the AR pipeline itself is robust.',
    fontsize=13, color=GRAY)

shape = s.shapes.add_shape(1, Inches(6.9), Inches(3.6), Inches(6.0), Inches(1.5))
shape.fill.solid(); shape.fill.fore_color.rgb = CARD
shape.line.color.rgb = RGBColor(0x2d, 0x2d, 0x2f); shape.line.width = Pt(1)
box(s, 7.0, 3.65, 2.5, 0.3, text='● LIVE STATS', fontsize=9, bold=True, color=GRAY)
box(s, 7.0, 4.0,  1.5, 0.5, text='29.8', fontsize=28, bold=True, color=GREEN)
box(s, 7.0, 4.5,  1.5, 0.3, text='FPS live', fontsize=11, color=GRAY)
box(s, 9.2, 4.0,  1.5, 0.5, text='0.60', fontsize=28, bold=True, color=BLUE)
box(s, 9.2, 4.5,  2.5, 0.3, text='Confidence threshold', fontsize=11, color=GRAY)

# ── SLIDE 9 — Limitations ────────────────────────────────────────────────────
s = add_slide(); bg(s)
eyebrow(s, 'LIMITATIONS & NEXT STEPS')
title(s, 'Honest about the gaps')

box(s, 0.5, 1.85, 5.8, 0.3, text='CURRENT LIMITATIONS',
    fontsize=10, bold=True, color=RED)
limits = [
    ('Fixed center ROI', 'Only detects objects at center of frame'),
    ('Fixed depth (0.5 m)', 'Box scale is wrong at very close or far distances'),
    ('No background class', 'Always outputs a class even when nothing is there'),
    ('Prediction flicker', 'No temporal smoothing between frames'),
]
for i, (title_txt, desc_txt) in enumerate(limits):
    top = 2.2 + i * 1.1
    circ = s.shapes.add_shape(9, Inches(0.5), Inches(top+0.05), Inches(0.38), Inches(0.38))
    circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x3a, 0x10, 0x10)
    circ.line.fill.background()
    box(s, 0.5, top+0.02, 0.38, 0.38, text='!',
        fontsize=13, bold=True, color=RED, align=PP_ALIGN.CENTER)
    box(s, 1.05, top,      5.2, 0.3, text=title_txt, fontsize=14, bold=True, color=WHITE)
    box(s, 1.05, top+0.32, 5.2, 0.4, text=desc_txt,  fontsize=12, color=GRAY)

box(s, 7.0, 1.85, 5.8, 0.3, text='EASY WINS NEXT',
    fontsize=10, bold=True, color=BLUE)
nexts = [
    'Collect webcam images for the 5 web-only classes → push accuracy above 90%',
    'Add a 5-frame history buffer to smooth out prediction flickering',
    'Add a background class so the system can say "I don\'t know"',
    'Replace the fixed ROI with a sliding-window or YOLO-based detector',
]
for i, txt in enumerate(nexts):
    top = 2.2 + i * 1.1
    circ = s.shapes.add_shape(9, Inches(7.0), Inches(top+0.05), Inches(0.38), Inches(0.38))
    circ.fill.solid(); circ.fill.fore_color.rgb = BLUE
    circ.line.fill.background()
    box(s, 7.0, top+0.02, 0.38, 0.38, text='→',
        fontsize=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    box(s, 7.55, top, 5.2, 0.65, text=txt, fontsize=13, color=WHITE)

# ── SLIDE 10 — Summary ───────────────────────────────────────────────────────
s = add_slide(); bg(s)
eyebrow(s, 'SUMMARY')
box(s, 0.5, 0.9, 12.33, 1.4,
    text='Any object the network knows\nbecomes its own AR anchor.',
    fontsize=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

stats = [
    ('10',     'object classes',        BLUE),
    ('84.1%',  'test accuracy',         GREEN),
    ('30 fps', 'live inference',        WHITE),
    ('0',      'printed markers needed',YELLOW),
]
for i, (val, lbl, col) in enumerate(stats):
    left = 0.5 + i * 3.2
    shape = s.shapes.add_shape(1, Inches(left), Inches(2.7), Inches(3.0), Inches(1.6))
    shape.fill.solid(); shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = RGBColor(0x2d, 0x2d, 0x2f); shape.line.width = Pt(1)
    box(s, left, 2.85, 3.0, 0.7, text=val,
        fontsize=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    box(s, left, 3.6,  3.0, 0.4, text=lbl,
        fontsize=12, color=GRAY, align=PP_ALIGN.CENTER)

box(s, 0, 5.6, 13.33, 0.4,
    text='Sangeeth Deleep Menon  ·  CS5330 Spring 2026  ·  Northeastern University',
    fontsize=12, color=GRAY, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out = os.path.join(BASE, 'presentation.pptx')
prs.save(out)
print('Saved:', out)
